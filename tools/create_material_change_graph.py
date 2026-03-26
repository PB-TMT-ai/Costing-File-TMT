#!/usr/bin/env python3
"""Generate a JSW One branded PowerPoint with month-on-month change bar charts.

Uses matplotlib for chart rendering (as images) and python-pptx for slide
layout with JSW One branding (logo, divider, colors, Calibri font).

Usage:
    python tools/create_material_change_graph.py [--input <path>] [--output <path>]

Exit codes:
    0 -- Success
    1 -- Error
    2 -- File not found or invalid input
"""

import argparse
import os
import shutil
import sys
import tempfile
from datetime import datetime
from pathlib import Path

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.ticker as ticker
except ImportError:
    print("ERROR: matplotlib not installed. Run: pip install matplotlib", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

import openpyxl

# ═══════════════════════════════════════════════════════════════
# JSW ONE BRAND CONSTANTS
# ═══════════════════════════════════════════════════════════════
JSW_BLUE = RGBColor(0x21, 0x33, 0x66)
JSW_RED = RGBColor(0xEA, 0x21, 0x27)
JSW_GREY = RGBColor(0x7F, 0x7F, 0x7F)
JSW_LTGREY = RGBColor(0xF2, 0xF2, 0xF2)
JSW_BORDER = RGBColor(0xCC, 0xCC, 0xCC)
JSW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
JSW_BLACK = RGBColor(0x00, 0x00, 0x00)

# Matplotlib hex colors matching JSW One palette
MPL_BLUE = "#213366"
MPL_GREY = "#7F7F7F"
MPL_LTGREY = "#F2F2F2"
MPL_BLUE_LIGHT = "#4A6A9E"

FONT_NAME = "Calibri"

# Slide dimensions (LAYOUT_WIDE)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Grid constants
GRID_L = 0.50
GRID_W = 12.33
GRID_TOP = 0.95
GRID_BOTTOM = 6.80
GRID_H = 5.85
FOOTER_Y = 7.05

# Logo path
LOGO_PATH = None
for p in [
    Path(__file__).parent.parent / ".claude" / "skills" / "jsw-one-pptx" / "assets" / "JSW_Logo_Clean.png",
    Path(__file__).parent.parent / ".claude" / "skills" / "jsw-one-pptx" / "assets" / "JSW_Logo_Final.png",
]:
    if p.exists():
        LOGO_PATH = str(p)
        break

# ═══════════════════════════════════════════════════════════════
# GRAPH SPECIFICATIONS
# ═══════════════════════════════════════════════════════════════
GRAPH_SPECS = [
    {"item": "Pallet DRI",         "row": 3,  "markets": ["Raipur"],        "unit": "INR/MT"},
    {"item": "Pig Iron",           "row": 4,  "markets": ["Raipur"],        "unit": "INR/MT"},
    {"item": "Scrap",              "row": 5,  "markets": ["Raipur", "NCR"], "unit": "INR/MT"},
    {"item": "Silico Manganese",   "row": 6,  "markets": ["Raipur"],        "unit": "INR/kg"},
    {"item": "Iron Ore DRI",       "row": 7,  "markets": ["Raipur"],        "unit": "INR/MT"},
    {"item": "Nett Margin Billet (Raipur)", "row": 10, "markets": ["Raipur"], "market_key": "Raipur", "unit": "INR/MT"},
    {"item": "Nett Margin Billet (NCR)",    "row": 10, "markets": ["NCR"],    "market_key": "NCR",    "unit": "INR/MT"},
    {"item": "Margin TMT (Raipur)",         "row": 12, "markets": ["Raipur"], "market_key": "Raipur", "unit": "INR/MT"},
    {"item": "Margin TMT (NCR)",            "row": 12, "markets": ["NCR"],    "market_key": "NCR",    "unit": "INR/MT"},
]

# ═══════════════════════════════════════════════════════════════
# DATA LOADING
# ═══════════════════════════════════════════════════════════════
def load_change_log(path):
    wb = openpyxl.load_workbook(path)
    ws = wb["Change Log"]
    dates, date_cols = [], []
    for col in range(2, ws.max_column + 1, 2):
        val = ws.cell(row=1, column=col).value
        if val is not None:
            dates.append(str(val))
            date_cols.append(col)
    data, seen_rows = {}, {}
    for spec in GRAPH_SPECS:
        row = spec["row"]
        if row in seen_rows:
            data[spec["item"]] = seen_rows[row]
            continue
        r_vals, n_vals = [], []
        for col in date_cols:
            r_vals.append(_parse_value(ws.cell(row=row, column=col).value))
            n_vals.append(_parse_value(ws.cell(row=row, column=col + 1).value))
        row_data = {"Raipur": r_vals, "NCR": n_vals}
        seen_rows[row] = row_data
        data[spec["item"]] = row_data
    wb.close()
    return dates, data


def _parse_value(val):
    if val is None or str(val).strip() in ("-", ""):
        return None
    try:
        return float(val)
    except (ValueError, TypeError):
        return None


def compute_changes(values, dates):
    changes, labels = [], []
    for i in range(1, len(values)):
        if values[i] is not None and values[i - 1] is not None:
            changes.append(values[i] - values[i - 1])
        else:
            changes.append(None)
        labels.append(_format_date_label(dates[i]))
    valid = [c for c in changes if c is not None]
    avg = sum(valid) / len(valid) if valid else 0
    return changes, labels, avg


def _format_date_label(date_str):
    try:
        dt = datetime.strptime(date_str, "%Y-%m-%d")
        return dt.strftime("%b'%y")
    except ValueError:
        return date_str[:7]


# ═══════════════════════════════════════════════════════════════
# CHART IMAGE GENERATION (matplotlib with JSW One colors)
# ═══════════════════════════════════════════════════════════════
def create_chart_image(spec, dates, data, tmp_dir):
    item = spec["item"]
    markets = spec["markets"]
    unit = spec["unit"]
    is_dual = len(markets) == 2 and "market_key" not in spec

    fig, ax = plt.subplots(figsize=(12.33, 5.0), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    if is_dual:
        _plot_dual_market(ax, spec, dates, data)
    else:
        _plot_single_market(ax, spec, dates, data)

    ax.set_ylabel(f"Change ({unit})", fontsize=11, color=MPL_GREY)
    ax.axhline(y=0, color=MPL_GREY, linewidth=0.5)
    ax.grid(axis="y", alpha=0.25, linestyle="--", color="#CCCCCC")
    ax.tick_params(axis="x", rotation=45, labelsize=10)
    ax.tick_params(axis="y", labelsize=10, colors=MPL_GREY)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CCCCCC")
    ax.spines["bottom"].set_color("#CCCCCC")

    plt.tight_layout()
    path = os.path.join(tmp_dir, f"{item.replace(' ', '_')}.png")
    fig.savefig(path, bbox_inches="tight", facecolor="white", dpi=150)
    plt.close(fig)
    return path


def _plot_single_market(ax, spec, dates, data):
    item = spec["item"]
    market_key = spec.get("market_key", "Raipur")
    values = data[item][market_key]
    changes, labels, avg = compute_changes(values, dates)

    x = range(len(labels))
    plot_changes = [c if c is not None else 0 for c in changes]

    bars = ax.bar(x, plot_changes, color=MPL_BLUE, width=0.6,
                  edgecolor="white", linewidth=0.5)
    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=10, fontfamily="sans-serif")

    _add_bar_labels(ax, bars, plot_changes, spec)

    # Average line
    ax.axhline(y=avg, color=MPL_BLUE_LIGHT, linewidth=1.5, linestyle="--", alpha=0.8)
    ax.annotate(
        f"Avg: {_fmt_val(avg, spec)}",
        xy=(len(labels) - 1, avg), fontsize=10, fontweight="bold",
        color=MPL_BLUE, ha="right",
        va="bottom" if avg >= 0 else "top",
    )


def _plot_dual_market(ax, spec, dates, data):
    item = spec["item"]
    r_values = data[item]["Raipur"]
    n_values = data[item]["NCR"]
    r_changes, labels, r_avg = compute_changes(r_values, dates)
    n_changes, _, n_avg = compute_changes(n_values, dates)

    x = list(range(len(labels)))
    width = 0.35
    r_plot = [c if c is not None else 0 for c in r_changes]
    n_plot = [c if c is not None else 0 for c in n_changes]

    x_r = [xi - width / 2 for xi in x]
    x_n = [xi + width / 2 for xi in x]

    bars_r = ax.bar(x_r, r_plot, width=width, color=MPL_BLUE, label="Raipur",
                    edgecolor="white", linewidth=0.5)
    bars_n = ax.bar(x_n, n_plot, width=width, color=MPL_GREY, label="NCR",
                    edgecolor="white", linewidth=0.5)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=10, fontfamily="sans-serif")

    _add_bar_labels(ax, bars_r, r_plot, spec, fontsize=7)
    _add_bar_labels(ax, bars_n, n_plot, spec, fontsize=7)

    ax.axhline(y=r_avg, color=MPL_BLUE, linewidth=1.5, linestyle="--", alpha=0.6)
    ax.axhline(y=n_avg, color=MPL_GREY, linewidth=1.5, linestyle="--", alpha=0.6)

    ax.annotate(f"Raipur avg: {_fmt_val(r_avg, spec)}", xy=(0, r_avg),
                fontsize=9, fontweight="bold", color=MPL_BLUE,
                va="bottom" if r_avg >= 0 else "top")
    ax.annotate(f"NCR avg: {_fmt_val(n_avg, spec)}", xy=(len(labels) - 1, n_avg),
                fontsize=9, fontweight="bold", color=MPL_GREY, ha="right",
                va="bottom" if n_avg >= 0 else "top")

    ax.legend(loc="upper left", fontsize=10, framealpha=0.9)


def _add_bar_labels(ax, bars, values, spec, fontsize=8):
    for bar, val in zip(bars, values):
        if val == 0:
            continue
        y = bar.get_height()
        va = "bottom" if y >= 0 else "top"
        ax.text(bar.get_x() + bar.get_width() / 2, y, _fmt_val(val, spec),
                ha="center", va=va, fontsize=fontsize, fontweight="bold",
                color=MPL_BLUE if val >= 0 else "#8B0000")


def _fmt_val(val, spec):
    if "Silico Manganese" in spec["item"]:
        return f"{val:+.1f}"
    if abs(val) >= 1:
        return f"{val:+,.0f}"
    return f"{val:+.2f}"


# ═══════════════════════════════════════════════════════════════
# JSW ONE BRANDED SLIDE HELPERS (python-pptx)
# ═══════════════════════════════════════════════════════════════
def _add_divider(slide, y_inches):
    """Blue + red divider bar."""
    # Blue segment (left ~66%)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y_inches),
                                   Inches(8.80), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = JSW_BLUE
    shape.line.fill.background()
    # Red segment (right ~36%, overlapping)
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.40), Inches(y_inches),
                                    Inches(4.93), Inches(0.05))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = JSW_RED
    shape2.line.fill.background()


def _add_logo(slide, slide_type="content"):
    if not LOGO_PATH:
        return
    if slide_type == "title":
        slide.shapes.add_picture(LOGO_PATH, Inches(10.30), Inches(2.10),
                                 Inches(2.40), Inches(0.79))
    else:
        slide.shapes.add_picture(LOGO_PATH, Inches(10.80), Inches(0.02),
                                 Inches(2.10), Inches(0.69))


def _add_page_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(0.30), Inches(FOOTER_Y),
                                     Inches(0.50), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.name = FONT_NAME
    p.font.color.rgb = JSW_GREY


def _add_source(slide, text):
    txBox = slide.shapes.add_textbox(Inches(6.50), Inches(FOOTER_Y),
                                     Inches(6.30), Inches(0.35))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.name = FONT_NAME
    p.font.color.rgb = JSW_GREY
    p.font.italic = True
    p.alignment = PP_ALIGN.RIGHT


def _add_text(slide, x, y, w, h, text, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.font.color.rgb = color or JSW_BLACK
    p.alignment = align
    return txBox


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════
def build_presentation(chart_paths, specs, dates, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 1
    _add_title_slide(prs, dates, page)
    page += 1

    for chart_path, spec in zip(chart_paths, specs):
        _add_chart_slide(prs, chart_path, spec, page)
        page += 1

    prs.save(output_path)
    print(f"Presentation saved: {output_path}", file=sys.stderr)


def _add_title_slide(prs, dates, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    _add_text(slide, 0.60, 1.90, 9.40, 1.00,
              "Material cost change analysis",
              size=28, bold=True, color=JSW_BLUE)
    _add_divider(slide, 2.95)
    _add_logo(slide, "title")
    _add_text(slide, 0.64, 3.12, 9.40, 0.45,
              "Month-on-month absolute change (INR) for key raw materials and margins",
              size=12, color=JSW_GREY)
    _add_text(slide, 0.64, 3.52, 9.40, 0.35,
              f"Period: {dates[0]} to {dates[-1]}  |  {len(dates)} data points",
              size=12, color=JSW_GREY)
    _add_text(slide, 0.64, 4.00, 9.40, 0.35,
              f"Generated: {datetime.now().strftime('%Y-%m-%d')}",
              size=12, color=JSW_GREY)
    _add_page_number(slide, page_num)


def _add_chart_slide(prs, chart_path, spec, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Slide heading
    if "market_key" in spec:
        title = f"{spec['item']} — {spec['unit']}"
    else:
        title = f"{spec['item']} ({' & '.join(spec['markets'])}) — {spec['unit']}"

    _add_text(slide, GRID_L, 0.10, 10.05, 0.58,
              title, size=20, bold=True, color=JSW_BLUE)

    _add_divider(slide, 0.75)
    _add_logo(slide, "content")

    # Chart image — fills content area
    slide.shapes.add_picture(chart_path,
                             Inches(GRID_L), Inches(GRID_TOP),
                             Inches(GRID_W), Inches(GRID_H))

    _add_page_number(slide, page_num)
    _add_source(slide, "Source: Costing change log")


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════
def main():
    parser = argparse.ArgumentParser(description="Generate JSW One branded material change graph PowerPoint")
    parser.add_argument("--input", "-i", default="output/change_log.xlsx",
                        help="Path to change_log.xlsx")
    parser.add_argument("--output", "-o", default="output/material_change_graphs.pptx",
                        help="Output PPTX path")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"ERROR: Input file not found: {args.input}", file=sys.stderr)
        sys.exit(2)

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)

    print(f"Loading change log: {args.input}", file=sys.stderr)
    dates, data = load_change_log(args.input)
    print(f"Found {len(dates)} dates: {dates[0]} to {dates[-1]}", file=sys.stderr)

    tmp_dir = tempfile.mkdtemp(prefix="material_charts_")
    chart_paths = []
    try:
        for spec in GRAPH_SPECS:
            print(f"  Generating chart: {spec['item']}...", file=sys.stderr)
            path = create_chart_image(spec, dates, data, tmp_dir)
            chart_paths.append(path)

        build_presentation(chart_paths, GRAPH_SPECS, dates, args.output)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    print(f"Done! {len(GRAPH_SPECS)} charts generated in {args.output}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
